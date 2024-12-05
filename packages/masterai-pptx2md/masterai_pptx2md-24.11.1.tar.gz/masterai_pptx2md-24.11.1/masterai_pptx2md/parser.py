import io
import logging
from tqdm import tqdm
from operator import attrgetter
from typing import List, Union, Optional, IO

from pptx.slide import Slide
from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import Shape
from pptx.presentation import Presentation as PresentationClass
from pptx.api import Presentation as PresentationFunc
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.text.text import _Paragraph as Paragraph, _Run as Run

from masterai_pptx2md.models import Config
from masterai_pptx2md.outputter import MarkDownOutPutter
from masterai_pptx2md.utils import (
    is_title,
    is_text_block,
    is_list_block,
    is_accent,
    is_strong,
    image_to_base64,
)

logger = logging.getLogger(__name__)


class Parse:

    def __init__(self, config: Config, out_putter: MarkDownOutPutter) -> None:
        self.config = config
        self.out_putter = out_putter

    def get_formatted_text(self, para: Paragraph) -> str:
        res: str = ""
        for run in para.runs:
            run: Run
            text: str = run.text
            if text == "":
                continue
            if not self.config.disable_escaping:
                text: str = self.out_putter.get_escaped(text)
            try:
                if run.hyperlink.address:
                    text: str = self.out_putter.get_hyperlink(
                        text, run.hyperlink.address
                    )
            except:
                text = self.out_putter.get_hyperlink(
                    text, "error:ppt-link-parsing-issue"
                )
            if is_accent(run.font):
                text: str = self.out_putter.get_accent(text)
            elif is_strong(run.font):
                text: str = self.out_putter.get_strong(text)
            if not self.config.disable_color:
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    text: str = self.out_putter.get_colored(text, run.font.color.rgb)
            res += text
        return res.strip()

    def process_title(
        self, shape: Union[SlidePlaceholder, Shape], slide_idx: int
    ) -> None:
        text: str = shape.text_frame.text.strip()
        self.out_putter.put_title(text, 1)

    def process_text_block(self, shape: Union[SlidePlaceholder, Shape], _: int) -> None:
        if is_list_block(shape):
            # generate list block
            for para in shape.text_frame.paragraphs:
                para: Paragraph
                if para.text.strip() == "":
                    continue
                text: str = self.get_formatted_text(para)
                self.out_putter.put_list(text, para.level)
            self.out_putter.write("\n")
        else:
            # generate paragraph block
            for para in shape.text_frame.paragraphs:
                para: Paragraph
                if para.text.strip() == "":
                    continue
                text: str = self.get_formatted_text(para)
                self.out_putter.put_para(text)

    def process_notes(self, text: str, _: int) -> None:
        self.out_putter.put_para("---")
        self.out_putter.put_para(text)

    def process_picture(self, shape: Picture, slide_idx: int) -> None:
        if self.config.disable_image:
            return
        pic_ext: str = shape.image.ext  # bmp gif jpg png tiff wmf
        image_bytes: bytes = shape.image.blob
        base64_image: Optional[str] = image_to_base64(image_bytes, pic_ext)
        if not base64_image:
            return
        self.out_putter.put_image(base64_image, self.config.max_img_width)

    def process_table(self, shape: GraphicFrame, _: int) -> None:
        table: List[List[str]] = [
            [cell.text for cell in row.cells] for row in shape.table.rows
        ]
        if len(table) > 0:
            self.out_putter.put_table(table)

    def ungroup_shapes(
        self, shapes: SlideShapes
    ) -> List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]]:
        res: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = []
        for shape in shapes:  # type: ignore
            shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    res.extend(self.ungroup_shapes(shape.shapes))  # type: ignore
                else:
                    res.append(shape)
            except Exception as e:
                logger.error(f"failed to load shape {shape}, skipped. error: {e}")
        return res

    def process_shapes(
        self,
        current_shapes: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]],
        slide_id: int,
    ) -> None:
        for shape in current_shapes:
            shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]
            if is_title(shape):
                self.process_title(shape, slide_id + 1)  # type: ignore
            elif is_text_block(shape):
                self.process_text_block(shape, slide_id + 1)  # type: ignore
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    self.process_picture(shape, slide_id + 1)  # type: ignore
                except AttributeError as e:
                    logger.error(f"Failed to process picture, skipped: {e}")
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self.process_table(shape, slide_id + 1)  # type: ignore
            else:
                try:
                    ph = shape.placeholder_format
                    if (
                        ph.type == PP_PLACEHOLDER.OBJECT
                        and hasattr(shape, "image")
                        and getattr(shape, "image")
                    ):
                        self.process_picture(shape, slide_id + 1)  # type: ignore
                except:
                    pass

    def parse(self, pptx_content: bytes) -> str:
        pptx: IO[bytes] = io.BytesIO(pptx_content)
        prs: PresentationClass = PresentationFunc(pptx)
        for idx, slide in enumerate(tqdm(prs.slides, desc="Converting slides")):
            idx: int
            slide: Slide
            shapes: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = []
            try:
                shapes_with_none: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = self.ungroup_shapes(slide.shapes)
                shapes: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = (
                    sorted(shapes_with_none, key=lambda x: (getattr(x, "top", 0) or 0, getattr(x, "left", 0) or 0))
                )
            except:
                logger.error(
                    "Bad shapes encountered in this slide. Please check or move them and try again."
                )
                logger.error("shapes:")
                try:
                    for sp in slide.shapes:
                        logger.error(sp.shape_type)
                        logger.error(sp.top, sp.left, sp.width, sp.height)
                except:
                    logger.error("failed to print all bad shapes.")

            self.process_shapes(shapes, idx + 1)

            if not self.config.disable_notes and slide.has_notes_slide:
                if not slide.notes_slide or not slide.notes_slide.notes_text_frame:
                    continue
                text: str = slide.notes_slide.notes_text_frame.text  # type: ignore
                if text:
                    self.process_notes(text, idx + 1)
            if idx < len(prs.slides) - 1 and self.config.enable_slides:
                self.out_putter.put_para("\n---\n")

        return self.out_putter.read()
