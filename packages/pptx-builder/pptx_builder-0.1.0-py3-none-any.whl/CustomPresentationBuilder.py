from pptx import Presentation
from pptx.util import Inches, Pt, Mm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class CustomPresentationBuilder:
    def __init__(self, slides_content_list):
        self.slides_content_list = slides_content_list
        self.prs = Presentation()
        # Set slide size to 16:9 aspect ratio
        self.prs.slide_width = Inches(13.3333)  # Width for 16:9
        self.prs.slide_height = Inches(7.5)  # Height for 16:9
        self.layout_mapping = self._create_layout_mapping()
        self.default_fonts = {
            'title': {'name': 'Arial', 'size': Pt(26.7), 'bold': False},
            'body': {'name': 'Arial', 'size': Pt(14), 'bold': False},
            'header': {'name': 'ＭＳ Ｐゴシック (本文)', 'size': Pt(11), 'bold': True},
            'cell': {'name': 'ＭＳ Ｐゴシック (本文)', 'size': Pt(11), 'bold': False},
        }

    def _create_layout_mapping(self):
        layout_mapping = {}
        for layout in self.prs.slide_layouts:
            name = layout.name.upper().replace(' ', '_')
            layout_mapping[name] = layout
        return layout_mapping

    def _set_text_properties(self, text_frame, text, font_props):
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        font = p.runs[0].font
        font.name = font_props.get('name', 'Arial')
        font.size = font_props.get('size', Pt(14))
        font.bold = font_props.get('bold', False)

    def _get_placeholder_by_idx(self, slide, idx):
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == idx:
                return placeholder
        return None

    def create_presentation(self, output_file):
        for slide_content in self.slides_content_list:
            layout_name = slide_content.get('layout', 'TITLE_AND_CONTENT').upper()
            layout = self.layout_mapping.get(layout_name, self.prs.slide_layouts[1])
            slide = self.prs.slides.add_slide(layout)

            # Set title
            title_placeholder = slide.shapes.title
            title_font_props = slide_content.get('title_font', self.default_fonts['title'])
            self._set_text_properties(title_placeholder.text_frame, slide_content.get('title', ''), title_font_props)

            # Handle different layouts
            if 'columns' in slide_content:
                columns = slide_content['columns']
                body_font_props = slide_content.get('body_font', self.default_fonts['body'])
                if layout_name == 'TWO_CONTENT' and len(columns) == 2:
                    left_placeholder = self._get_placeholder_by_idx(slide, 1)
                    right_placeholder = self._get_placeholder_by_idx(slide, 2)
                    if left_placeholder and right_placeholder:
                        self._set_text_properties(left_placeholder.text_frame, columns[0], body_font_props)
                        self._set_text_properties(right_placeholder.text_frame, columns[1], body_font_props)
                    else:
                        print("Error: Expected placeholders not found in slide layout.")
                else:
                    # Manually create columns
                    shapes = slide.shapes
                    num_columns = len(columns)
                    column_width = self.prs.slide_width / num_columns
                    height = self.prs.slide_height - Inches(1.5)  # Adjust as needed
                    for idx, column_text in enumerate(columns):
                        left = column_width * idx
                        top = Inches(1.5)
                        txBox = shapes.add_textbox(left, top, column_width, height)
                        self._set_text_properties(txBox.text_frame, column_text, body_font_props)
            if 'text' in slide_content:
                # Single body text
                body_placeholder = self._get_placeholder_by_idx(slide, 1)
                body_font_props = slide_content.get('body_font', self.default_fonts['body'])
                if body_placeholder:
                    self._set_text_properties(body_placeholder.text_frame, slide_content['text'], body_font_props)
                else:
                    # Create a new textbox
                    shapes = slide.shapes
                    left = Inches(1)
                    top = Inches(2)
                    width = self.prs.slide_width - Inches(2)
                    height = self.prs.slide_height - Inches(3)
                    txBox = shapes.add_textbox(left, top, width, height)
                    self._set_text_properties(txBox.text_frame, slide_content['text'], body_font_props)
            if 'grid' in slide_content:
                # Grid of images
                grid = slide_content['grid']
                self._add_image_grid(slide, grid)
            if 'table' in slide_content:
                # Add table from DataFrame
                table_content = slide_content['table']
                df = table_content['dataframe']
                position = table_content['position']
                left_mm = position.get('left_mm', 10)
                top_mm = position.get('top_mm', 50)
                width_mm = position.get('width_mm', 200)
                height_mm = position.get('height_mm', 100)
                header_font_props = table_content.get('header_font', self.default_fonts['header'])
                cell_font_props = table_content.get('cell_font', self.default_fonts['cell'])
                self._add_table(slide, df, left_mm, top_mm, width_mm, height_mm, header_font_props, cell_font_props)
            else:
                # No additional content
                pass

        # Save the presentation
        self.prs.save(output_file)

    def _add_image_grid(self, slide, grid):
        # [Your existing _add_image_grid code remains unchanged]
        pass  # Placeholder for the existing code

    def _add_table(self, slide, df, left_mm, top_mm, width_mm, height_mm, header_font_props, cell_font_props):
        left = Mm(left_mm)
        top = Mm(top_mm)
        width = Mm(width_mm)
        height = Mm(height_mm)

        rows = df.shape[0] + 1  # Add one for the header row
        cols = df.shape[1]

        # Add table shape
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths (equal widths for simplicity)
        for col_idx in range(cols):
            table.columns[col_idx].width = int(width / cols)

        # Write the header row
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    font = run.font
                    font.name = header_font_props.get('name', 'ＭＳ Ｐゴシック (本文)')
                    font.size = header_font_props.get('size', Pt(11))
                    font.bold = header_font_props.get('bold', True)

        # Write the data rows
        for row_idx, (index, row) in enumerate(df.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        font = run.font
                        font.name = cell_font_props.get('name', 'ＭＳ Ｐゴシック (本文)')
                        font.size = cell_font_props.get('size', Pt(11))
                        font.bold = cell_font_props.get('bold', False)
