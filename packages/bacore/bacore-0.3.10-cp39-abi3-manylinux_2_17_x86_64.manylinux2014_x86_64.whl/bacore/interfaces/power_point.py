"""Power point interface."""

from bacore.domain.measurements import Time
from dataclasses import dataclass
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Inches, Length, Pt
from sqlmodel import SQLModel, Field
from typing import ClassVar

TODAY = Time().today_s


@dataclass(frozen=True)
class PowerPoint:
    """PowerPoint class with pptx.Presentation as `prs` attribute.

    Attributes:
        `prs`: Contains the pptx.Presentation object.

    Methods:
        `add_slide`: Add a slide from default template using template index and having an optional title.
    """

    prs = Presentation()

    background_layer: ClassVar[int] = 2
    widescreen_width: ClassVar[Inches] = Inches(13.33)
    widescreen_height: ClassVar[Inches] = Inches(7.5)

    def add_slide(self, layout_index: int, title_text: str | None = None) -> Slide:
        """Add a PowerPoint slide with an optional title text.

        Returns:
            slide object
        """
        if layout_index >= len(self.prs.slide_layouts):
            raise ValueError(f"Layout index '{layout_index}' out of range.")

        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        if title_text:
            title = slide.shapes.title
            title.text = title_text
        return slide

    @staticmethod
    def add_background_image(
        slide: Slide,
        image_file: str,
        left: Length = 0,
        top: Length = 0,
        width: Length | None = None,
        height: Length | None = None,
        move_to_background: bool = True,
    ) -> Picture:
        """Add a background image to a slide.

        Parameters:
            slide: The slide to which the image will be added.
            image_file: Path to the image file.
            left: The left position of the image.
            top: The top position of the image.
            width: The width of the image.
            height: The height of the image.
            move_to_background: If `True`, moves the image to the back of the stack.
        """
        background_img = slide.shapes.add_picture(image_file, left, top, width, height)
        if move_to_background:
            slide.shapes._spTree.remove(background_img._element)
            slide.shapes._spTree.insert(2, background_img._element)
        return background_img

    @classmethod
    def default_templates(cls, widescreen: bool | None = True) -> None:
        """Create a power point slide presentation using default templates."""
        prs = Presentation()

        if widescreen:
            prs.slide_width = cls.widescreen_width
            prs.slide_height = cls.widescreen_height

        for layout_index in range(len(prs.slide_layouts)):
            cls.add_slide(layout_index)

        prs.save("default_templates.pptx")


@dataclass(frozen=True)
class Placeholder:
    """A placeholder item inside of a slide."""

    slide: Slide
    id: int

    def add_bullets(self, bullets: list[tuple[str, int | None, int]]):
        """Add list of bullets to the slide.

        Each bullet is a tuple which consist of the text as a string, then an optional int for the bullet level, then another in for the font size.
        """
        requirement_body = self.slide.shapes.placeholders[self.id]
        requirement_body_tf = requirement_body.text_frame

        for bullet_text, bullet_level, font_size in bullets:
            requirement_body_tf_p = requirement_body_tf.add_paragraph()
            requirement_body_tf_p.text = bullet_text
            requirement_body_tf_p.font.size = Pt(font_size)
            if bullet_level is not None:
                requirement_body_tf_p.level = bullet_level


class TitleSlide(SQLModel, table=True):
    """Represents the title slide of a presentation."""

    id: int | None = Field(default=None, primary_key=True)
    title: str = Field(index=True)
    sub_title: str | None = Field(default=None)
    background_image: str | None = Field(default=None)
    date: str | None = Field(default=None)
    logo: str | None = Field(default=None)

    def create(self, ppt: PowerPoint) -> Slide:
        """Create a title slide."""
        slide = ppt.add_slide(0, self.title)
        if self.background_image:
            ppt.add_background_image(
                slide,
                self.background_image,
                width=PowerPoint.widescreen_width,
                height=Inches(7),
            )
        if self.sub_title:
            subtitle = slide.shapes.placeholders[1]
            subtitle.text = self.sub_title
            subtitle.text_frame.paragraphs[0].font.italic = True
            subtitle.text_frame.paragraphs[0].font.color.rbg = RGBColor(255, 0, 0)
        if self.date:
            date = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(7.1), width=Inches(1), height=Inches(0.3))
            date_tf = date.text_frame
            date_tf.text = self.date
            date_p = date_tf.paragraphs[0]
            date_run = date_p.runs[0]
            date_run.font.size = Pt(12)
            date_run.font.italic = True
            date_p.alignment = PP_ALIGN.CENTER
            date_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if self.logo:
            slide.shapes.add_picture(
                self.logo,
                left=Inches(11),
                top=Inches(7),
                width=Inches(1.5),
            )
        return slide
