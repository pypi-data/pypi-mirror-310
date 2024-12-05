"""Power point interface."""

from dataclasses import dataclass
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Inches, Length, Pt


def add_background_image(
    slide: Slide,
    image_file: str,
    left: Length = 0,
    top: Length = 0,
    width: Length | None = None,
    height: Length | None = None,
    move_to_background: bool = True,
) -> Picture:
    """Add a background image to a slide.

    Parameters:
        slide: The slide to which the image will be added.
        image_file: Path to the image file.
        left: The left position of the image.
        top: The top position of the image.
        width: The width of the image.
        height: The height of the image.
        move_to_background: If `True`, moves the image to the back of the stack.
    """
    background_img = slide.shapes.add_picture(image_file=image_file, left=left, top=top, width=width, height=height)
    if move_to_background:
        slide.shapes._spTree.remove(background_img._element)
        slide.shapes._spTree.insert(2, background_img._element)
    return background_img


class PowerPoint:
    """PowerPoint class with pptx.Presentation as `prs` attribute.

    Attributes:
        `prs`: Contains the pptx.Presentation object.

    Methods:
        `add_slide`: Add a slide from default template using template index and having an optional title.
    """

    widescreen_width = Inches(13.33)
    widescreen_height = Inches(7.5)

    def __init__(self) -> None:
        self.prs = Presentation()

    def add_slide(self, layout_index: int, title_text: str | None = None) -> Slide:
        """Add a PowerPoint slide with an optional title text.

        Returns:
            slide object
        """
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        if title_text:
            title = slide.shapes.title
            title.text = title_text
        return slide


@dataclass
class Placeholder:
    """A placeholder item inside of a slide."""

    slide: Slide
    id: int

    def add_bullets(self, bullets: list[tuple[str, int | None, int]]):
        """Add list of bullets to the slide.

        Each bullet is a tuple which consist of the text as a string, then an optional int for the bullet level, then another in for the font size.
        """
        requirement_body = self.slide.shapes.placeholders[self.id]
        requirement_body_tf = requirement_body.text_frame

        for bullet_text, bullet_level, font_size in bullets:
            requirement_body_tf_p = requirement_body_tf.add_paragraph()
            requirement_body_tf_p.text = bullet_text
            requirement_body_tf_p.font.size = Pt(font_size)
            if bullet_level:
                requirement_body_tf_p.level = bullet_level


def default_templates(widescreen_dimensions: bool | None = True) -> None:
    """Create a power point slide presentation using default templates."""
    ppt = PowerPoint()

    if widescreen_dimensions:
        ppt.prs.slide_width = PowerPoint.widescreen_width
        ppt.prs.slide_height = PowerPoint.widescreen_height

    for layout_index in range(11):
        ppt.add_slide(layout_index)

    ppt.prs.save("default_templates.pptx")
