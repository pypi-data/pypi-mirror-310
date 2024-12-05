import shutil
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from tqdm import tqdm

from gg.utils.common import ppt_temp01_path, ColorPrinter, delete_file, tmp_dir_path, get_path_str, ppt_temp02_path, \
    sanitize_filename


def detect_black_block_widths(image):
    # 1. 转为灰度图
    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    # 2. 二值化
    _, binary_image = cv2.threshold(gray_image, 50, 255, cv2.THRESH_BINARY_INV)
    # 3. 查找轮廓
    contours, _ = cv2.findContours(binary_image, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    # 4. 测量黑色方块宽度并筛选符合条件的宽度
    block_widths = []
    width_set = set()
    for contour in contours:
        x, y, w, h = cv2.boundingRect(contour)
        width_set.add(w)
        block_widths.append((x, w))
    # 获取最大宽度
    max_width = max(width_set)
    # 筛选出符合条件的方块
    ret = []
    for x, w in block_widths:
        if w < max_width - 1:  # 过滤掉小于最大宽度的方块
            continue
        ret.append((x, w))
    ret.sort()
    return ret


def crop_image_from_left(image, crop_width, fill_color=(255, 255, 255)):
    height, width, channels = image.shape
    crop_width = min(crop_width, width)
    result_image = np.full((height, width, channels), fill_color, dtype=np.uint8)
    result_image[:, :crop_width] = image[:, :crop_width]
    return result_image


def gen_split_k_image(image_path: Path):
    # 1. 读取图像
    with open(image_path, 'rb') as f:
        file_bytes = np.asarray(bytearray(f.read()), dtype=np.uint8)
        image = cv2.imdecode(file_bytes, cv2.IMREAD_COLOR)

    ret = detect_black_block_widths(image)
    img_path_list = [str(image_path)]
    for i, (r1, r2) in enumerate(tqdm(ret, desc="生成图片中...", unit="张")):
        # 计算截取范围
        tmp = crop_image_from_left(image, r1 + r2 + 1)

        # 保存结果
        output_filename = get_path_str(tmp_dir_path / f"{i + 1}-{r1}-{r2}.png")
        cv2.imwrite(output_filename, tmp)
        img_path_list.append(output_filename)

    ColorPrinter.print("开始生成ppt...")
    ppt_name = image_path.with_suffix(".pptx")
    create_ppt_from_template(img_path_list, ppt_name)
    ColorPrinter.print_green(f"{ppt_name}生成成功")

    for img_path in img_path_list:
        delete_file(img_path)


def images_to_ppt(dir_path: Path, ppt_name: str):
    ppt_name = Path(ppt_name)
    ppt_name = ppt_name.with_name(sanitize_filename(ppt_name.name)).with_suffix(".pptx")
    img_path_list = sorted(dir_path.glob("*.png"))
    create_ppt_from_template(img_path_list, ppt_name, ppt_temp02_path)
    ColorPrinter.print_green(f"{ppt_name}生成成功")
    for img_path in img_path_list:
        delete_file(img_path)


def create_ppt_from_template(image_list, output_pptx_path, template_path=None):
    # 复制模板文件到目标路径
    template_path = template_path or ppt_temp01_path
    shutil.copy(template_path, output_pptx_path)

    # 加载复制的模板文件
    prs = Presentation(output_pptx_path)

    # 获取幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    slide_count = len(prs.slides)
    # 遍历图片列表
    for idx, image_path in enumerate(image_list):
        image_path = str(image_path)
        if idx > slide_count - 1:
            continue
        # 使用模板的第一个布局
        slide = prs.slides[idx]

        # 打开图片并获取其尺寸
        img = Image.open(image_path)
        img_width, img_height = img.size

        # 计算图片的缩放比例
        ratio = min(slide_width / img_width, slide_height / img_height)
        new_width = int(img_width * ratio)
        new_height = int(img_height * ratio)

        # 计算图片居中位置
        left = (slide_width - new_width) // 2
        top = (slide_height - new_height) // 2

        # 添加图片到幻灯片
        slide.shapes.add_picture(image_path, left, top, new_width, new_height)

    # 保存修改后的文件
    prs.save(output_pptx_path)
    return output_pptx_path


if __name__ == '__main__':
    images_to_ppt(Path("D:\自动截图"), "测试.ppt")
